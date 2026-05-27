# -*- coding: utf-8 -*-
import streamlit as st
import pandas as pd
import google.generativeai as genai
from google.api_core import exceptions  # Rate Limit(429) エラーを確実に捕捉するため
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
import io
import configparser
import os
import signal
import re
import time  # リトライ待機（スリープ）処理のため
import requests  # Supabase REST API用
import json

# --- 1. APIキーの設定 (APIKEY.ini または クラウドのSecretsからハイブリッド取得) ---
# ※既存のAPI取得ロジックの構造・変数名を1行も崩さずに、クラウド安全対策を内包させています
def load_api_key():
    # 1. まずローカルの APIKEY.ini を探す
    config = configparser.ConfigParser()
    file_path = 'APIKEY.ini'
    if os.path.exists(file_path):
        try:
            config.read(file_path, encoding='utf-8-sig')
            return config.get('GEMINI', 'API_KEY')
        except:
            pass
            
    # 2. もしローカルにファイルがなければ、Streamlit Cloudの「Secrets」から安全に取得する
    try:
        if "GEMINI" in st.secrets and "API_KEY" in st.secrets["GEMINI"]:
            return st.secrets["GEMINI"]["API_KEY"]
        elif "API_KEY" in st.secrets:
            return st.secrets["API_KEY"]
    except:
        return None
        
    return None

INI_KEY = load_api_key()
EMBEDDED_API_KEY = INI_KEY

# --- Supabase 連携用関数（サイレント安全対策版） ---
def get_supabase_config():
    try:
        if "SUPABASE" in st.secrets:
            return st.secrets["SUPABASE"].get("URL"), st.secrets["SUPABASE"].get("KEY")
    except:
        pass
    return None, None

def save_log_to_supabase(username, persona_name, messages, report=None):
    url, key = get_supabase_config()
    if not url or not key:
        return False
    
    headers = {
        "apikey": key,
        "Authorization": f"Bearer {key}",
        "Content-Type": "application/json",
        "Prefer": "return=minimal"
    }
    
    data = {
        "username": username,
        "persona_name": persona_name,
        "messages": messages,
        "report": report
    }
    
    try:
        res = requests.post(f"{url}/rest/v1/roleplay_logs", headers=headers, json=data, timeout=3)
        return res.status_code in [200, 201]
    except Exception as e:
        return False

def fetch_logs_from_supabase():
    url, key = get_supabase_config()
    if not url or not key:
        return []
    
    headers = {
        "apikey": key,
        "Authorization": f"Bearer {key}",
        "Content-Type": "application/json"
    }
    
    try:
        res = requests.get(f"{url}/rest/v1/roleplay_logs?order=created_at.desc", headers=headers, timeout=3)
        if res.status_code == 200:
            return res.json()
    except Exception as e:
        pass
    return []

# --- 2. 各ファイル抽出関数 ---
def extract_from_docx(file):
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_from_pdf(file):
    reader = PdfReader(file)
    return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])

def extract_from_pptx(file):
    prs = Presentation(file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_from_excel(file):
    all_sheets = pd.read_excel(file, sheet_name=None)
    text_data = []
    for sheet_name, df in all_sheets.items():
        text_data.append(f"--- シート名: {sheet_name} ---\n{df.to_string(index=False)}")
    return "\n".join(text_data)

# --- Excel出力用の整形ロジック ---
def create_excel_download(text):
    output = io.BytesIO()
    lines = text.split('\n')
    table_data = []
    
    for line in lines:
        if '|' in line:
            cells = [c.strip() for c in line.split('|') if c.strip()]
            if cells and not all(c == '-' or c.startswith('---') for c in cells):
                table_data.append(cells)
                
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        if table_data:
            df = pd.DataFrame(table_data)
            df.to_excel(writer, index=False, header=False, sheet_name='評価結果')
            workbook = writer.book
            worksheet = writer.sheets['評価結果']
            border_fmt = workbook.add_format({'border': 1, 'text_wrap': True, 'valign': 'top'})
            for row_num in range(len(table_data)):
                worksheet.set_row(row_num, None, border_fmt)
            worksheet.set_column(0, 5, 30)
        else:
            df = pd.DataFrame([text.split('\n')])
            df.to_excel(writer, index=False, header=False, sheet_name='評価結果')
            
    return output.getvalue()

# --- 404エラーを回避しつつ、利用可能なモデル名を安全に取得する関数 ---
def get_safe_model_name(api_key):
    try:
        genai.configure(api_key=api_key)
        available_models = [m.name for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
        
        has_flash = any('gemini-1.5-flash' in m for m in available_models)
        target_raw = 'gemini-1.5-flash' if has_flash else (available_models[0] if available_models else 'gemini-1.5-flash')
        
        safe_name = target_raw.replace('models/', '')
        return safe_name
    except:
        return 'gemini-1.5-flash'

# --- ペルソナから最初の挨拶を動的に生成する関数 ---
def generate_initial_greeting(persona, api_key):
    target_model = get_safe_model_name(api_key)
    for attempt in range(5):
        try:
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel(target_model)
            
            prompt = f"""
あなたは以下の【顧客ペルソナ】として、銀行の窓口に手続きに来た本人です。
これからはじまる銀行員とのロールプレイに備えて、あなた（顧客）の側から最初にかける言葉（第一声）を1つだけ作成してください。

【顧客ペルソナ】
・氏名: {persona['name']}
・年齢: {persona['age']}
・職業: {persona['job']}
・家族構成: {persona['family']}
・来店目的: {persona['purpose']}
・性格: {persona['personality']}
・投資経験: {persona['experience']}

【作成ルール】
1. 基本的なベース文である「こんにちは。定期預金の満期の件で来たんですが、今の時代、普通に預けていても全然増えないですよね…」のニュアンスを守りつつ、上記の「来店目的」「性格」「年齢」に最も合致した自然な口調にアレンジしてください。
2. 1〜2文程度の短いセリフにしてください。
3. AIとしての挨拶や解説、メタ発言（例：「以下が最初の発言です」など）は一切含めず、顧客のセリフだけを出力してください。
"""
            response = model.generate_content(prompt)
            return response.text.strip()
        except (exceptions.ResourceExhausted, Exception) as e:
            error_msg = str(e)
            if "429" in error_msg or isinstance(e, exceptions.ResourceExhausted):
                time.sleep((attempt + 1) * 10)
                continue
            return "こんにちは。定期預金の満期の件で来たんですが、今の時代、普通に預けていても全然増えないですよね…"
    return "こんにちは。定期預金の満期の件で来たんですが、今の時代、普通に預けていても全然増えないですよね…"

# --- 3. AI回答生成ロジック (自動リトライ・履歴ウィンドウ削減版) ---
def get_ai_roleplay_response(messages, persona, product_docs, api_key):
    target_model = get_safe_model_name(api_key)
    recent_messages = [messages[0]] + messages[-5:] if len(messages) > 6 else messages
    
    for attempt in range(5):
        try:
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel(target_model)
            combined_docs = "\n\n".join(product_docs) if product_docs else "特になし"
            
            history_text = ""
            for m in recent_messages:
                role_label = "顧客(あなた)" if m["role"] == "assistant" else "銀行員(ユーザー)"
                history_text += f"{role_label}: {m['content']}\n"
                
            system_prompt = f"""
あなたは以下の【顧客ペルソナ】なりきり、銀行員（ユーザー）と対話をするロールプレイAIです。
設定された性格や知識レベル、感情をリアルに反映させて対話してください。

【顧客ペルソナ】
・氏名: {persona['name']}
・年齢: {persona['age']}
・職業: {persona['job']}
・家族構成: {persona['family']}
・来店目的: {persona['purpose']}
・性格: {persona['personality']}
・投資経験: {persona['experience']}

【参考：案内可能な金融商品情報】
{combined_docs}

【ロールプレイの絶対ルール】
1. ユーザー（銀行員）の返答に対して、設定された性格や知識レベルに基づいてリアルに会話を続けてください。
2. 一度に長文を話さないでください。実際の会話と同様に、1〜3文程度（最大でも100文字〜150文字程度）でテンポよく短く返答してください。
3. ユーザーが「専門用語（ポートフォリオ、コモディティ、信託報酬など）」を説明なしに使った場合、少し困惑した態度をとるか、質問し直してください。
4. 銀行員としてのマナーや、強引な勧誘、法令遵守（リスク説明の有無など）に問題があると感じた場合も、顧客の感情としてリアルに反応（例：「そんなにリスクがあるなら、やっぱりいいです」など）してください。
5. 会話が「商品の購入」「検討のため持ち帰り」「お断り」のいずれかの節目を迎えるまで、ロールプレイを継続してください。
6. 絶対に途中でAIとしてのメタな発言（「ロールプレイを終了しますか？」や「次の質問をどうぞ」など）をしないでください。あなたは100%このペルソナの人間です。

【これまでの会話履歴（※直近の重要な会話のみ抽出）】
{history_text}

上記のルールと履歴を元に、次につづく「顧客(あなた)」の短い発言を1つだけ生成してください。
AIとしての解説やメタ発言、余計な挨拶は一切含めないでください。
"""
            response = model.generate_content(system_prompt)
            return response.text
            
        except (exceptions.ResourceExhausted, Exception) as e:
            error_msg = str(e)
            if "429" in error_msg or isinstance(e, exceptions.ResourceExhausted):
                wait_time = (attempt + 1) * 10
                time.sleep(wait_time)
                continue
            return f"【システムエラー】詳細: {error_msg}"
            
    return "【混雑エラー】現在AIへのリクエストが連続しています。無料枠の制限を超過したため、1分ほど待ってから再度送信してください。"

# --- 4. 評価レポート生成ロジック (自動リトライ版) ---
def generate_evaluation_report(messages, persona, api_key):
    target_model = get_safe_model_name(api_key)
    
    for attempt in range(5):
        try:
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel(target_model)
            
            history_text = ""
            for m in messages:
                role_label = "顧客" if m["role"] == "assistant" else "銀行員(ユーザー)"
                history_text += f"{role_label}: {m['content']}\n"
                
            evaluation_prompt = f"""
    あなたは銀行のコンプライアンスおよび営業指導の専門家です。
    以下の【顧客ペルソナ】を相手に行われた、銀行員（ユーザー）のロールプレイ音声・チャットログを厳密に評価し、詳細な「評価レポート」を作成してください。
    
    【対象顧客ペルソナ】
    氏名: {persona['name']}, 年齢: {persona['age']}, 職業: {persona['job']}, 投資経験: {persona['experience']}
    
    【会話ログ】
    {history_text}
    
    【評価基準】
    1. 適切な案内ができているか (顧客のニーズ・来店目的に合致した提案か、ヒアリングは十分か)
    2. コンプライアンス上問題となる発言をしていないか (不確実な事項の断定的判断、リスクの不告知、強引な勧誘、適合性の原則違反がないか)
    3. 顧客本位の業務運営 (専門用語をわかりやすく噛み砕いて説明できているか)
    4. ビジネスマナー・コミュニケーション (適切な言葉遣い、傾聴の姿勢ができているか)
    5. 意向把握とニーズ深掘り (表面的な来店目的だけでなく、顧客の潜在的な資金ニーズやリスク許容度を適切に引き出せているか)
    
    【出力フォーマット】
    以下の構成で、マークダウンの「表形式」を活用して分かりやすく出力してください。
    テーブルのヘッダー（行頭）と区切り線の間に余計な改行やスペースを絶対に入れないでください。
    総合判定は（A: 優秀、B: 合格、C: 要指導）の3段階で評価してください。
    
    ### 📊 総合評価: [A / B / C]
    ### 📝 項目別詳細評価
    | 評価項目 | 評価点 (5点満点) | 具体的フィードバック・指摘事項 |
    | --- | --- | --- |
    | 適切な案内・提案 | | |
    | コンプライアンス遵守 | | |
    | 顧客本位（用語説明等） | | |
    | マナー・コミュニケーション | | |
    | 意向把握・ニーズ深掘り | | |
    
    ### 💡 良かった点
    （箇条書きで記入）
    
    ### ⚠️ 改善すべき点・指導事項
    （箇条書きで記入、特に具体的なコンプライアンス上のリスクやNGワードがあれば指摘）
    """
            response = model.generate_content(evaluation_prompt)
            
            # Streamlitでのマークダウン表崩れを防ぐための自動クリーニングロジック
            cleaned_text = response.text
            cleaned_text = re.sub(r'\|\s*[-–—]+\s*\|\s*[-–—]+\s*\|\s*[-–—]+\s*\|', '| --- | --- | --- |', cleaned_text)
            cleaned_text = cleaned_text.replace("項目別詳細評価\n\n|", "項目別詳細評価\n|")
            
            return cleaned_text
            
        except (exceptions.ResourceExhausted, Exception) as e:
            error_msg = str(e)
            if "429" in error_msg or isinstance(e, exceptions.ResourceExhausted):
                wait_time = (attempt + 1) * 15
                time.sleep(wait_time)
                continue
            return f"評価レポートの生成に失敗しました: {error_msg}"
            
    return "【混雑エラー】評価レポート生成リクエストが無料制限に達しました。1分ほど時間を空けて再度お試しください。"

# --- 5. 画面構築 (Streamlit UI) ---
st.set_page_config(page_title="銀行員向け 金融商品販売AIロールプレイ", layout="wide")

# --- ログイン機能の実装 ---
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False
if "username" not in st.session_state:
    st.session_state.username = ""

if not st.session_state.logged_in:
    st.title("🔐 ロールプレイシステム ログイン")
    with st.form("login_form"):
        input_user = st.text_input("ユーザー名")
        input_pass = st.text_input("パスワード", type="password")
        submit_login = st.form_submit_button("ログイン")
        
        if submit_login:
            if (input_user == "demo" and input_pass == "abic5980") or (input_user == "kanri" and input_pass == "abic5980"):
                st.session_state.logged_in = True
                st.session_state.username = input_user
                st.success(f"{input_user} としてログインしました。")
                st.rerun()
            else:
                st.error("ユーザー名またはパスワードが正しくありません。")
    st.stop()

# --- 管理者画面の表示ロジック（強制リフレッシュ st.rerun() 完備版） ---
if st.session_state.username == "kanri":
    st.title("🖥️ 管理者専用 ダッシュボード")
    st.markdown(f"ログイン中: {st.session_state.username} (管理者)")
    
    if st.button("🚪 ログアウト", key="admin_logout"):
        st.session_state.logged_in = False
        st.session_state.username = ""
        if "selected_index" in st.session_state: del st.session_state["selected_index"]
        if "prev_user" in st.session_state: del st.session_state["prev_user"]
        st.rerun()
        
    st.write("---")
    st.subheader("📋 担当者別の応答記録・評価結果一覧")
    
    # セッション記憶の初期化
    if "selected_index" not in st.session_state:
        st.session_state["selected_index"] = 0
    if "prev_user" not in st.session_state:
        st.session_state["prev_user"] = "すべて"
    
    with st.spinner("Supabaseから最新データを取得中..."):
        logs = fetch_logs_from_supabase()
        
    if not logs:
        st.info("保存されているロールプレイ記録はありません、またはSupabaseとの通信にタイムアウトしました。")
    else:
        log_list = []
        for l in logs:
            log_list.append({
                "日時": l.get("created_at"),
                "担当者": l.get("username"),
                "ペルソナ名": l.get("persona_name"),
                "対話回数": len(l.get("messages", [])),
                "評価レポート有無": "あり" if l.get("report") else "なし",
                "raw_data": l
            })
        df_base = pd.DataFrame(log_list)
        
        # 担当者でのフィルター変更時に関数で即時リフレッシュさせる
        def on_change_user_filter():
            st.session_state["selected_index"] = 0  # フィルタ変更時はインデックスを初期化
            
        users_list = ["すべて"] + list(df_base["担当者"].unique())
        selected_user = st.selectbox(
            "担当者でフィルター", 
            options=users_list,
            index=users_list.index(st.session_state["prev_user"]) if st.session_state["prev_user"] in users_list else 0,
            key="current_user_filter_box",
            on_change=lambda: st.session_state.update({"prev_user": st.session_state["current_user_filter_box"], "selected_index": 0})
        )
        
        if st.session_state["prev_user"] != "すべて":
            df_filtered = df_base[df_base["担当者"] == st.session_state["prev_user"]].reset_index(drop=True)
        else:
            df_filtered = df_base.reset_index(drop=True)
            
        # フィルター後のデータフレームを画面に表示
        st.dataframe(df_filtered.drop(columns=["raw_data"]), use_container_width=True)
        
        st.write("---")
        st.subheader("🔍 詳細ログ確認")
        
        if len(df_filtered) > 0:
            options_indices = list(range(len(df_filtered)))
            
            if st.session_state["selected_index"] >= len(df_filtered):
                st.session_state["selected_index"] = 0
                
            # 💡【重要】プルダウン選択が切り替わった瞬間に、セッションを上書きして st.rerun() で画面全体を再描画させる
            def on_change_selection():
                # 選択された値の番号を取得
                chosen_val = st.session_state["current_selection_box"]
                st.session_state["selected_index"] = chosen_val
                # 🔴 画面を強制再描画させて下部エリアへ即座に同期・反映させる
                st.rerun()
                
            selected_box_val = st.selectbox(
                "確認したいデータの番号を選択してください", 
                options=options_indices,
                index=st.session_state["selected_index"],
                format_func=lambda x: f"No.{x} : {df_filtered.iloc[x]['日時']} - {df_filtered.iloc[x]['担当者']} ({df_filtered.iloc[x]['ペルソナ名']})",
                key="current_selection_box",
                on_change=on_change_selection
            )
            
            # セッション状態に保存されたインデックス番号から確実にデータを追従抽出
            target_log = df_filtered.iloc[st.session_state["selected_index"]]["raw_data"]
            
            st.markdown("---")
            st.markdown(f"### 🎯 選択されたデータ詳細 (ペルソナ: **{target_log.get('persona_name')}** / 担当者: **{target_log.get('username')}**)")
            
            col_admin_chat, col_admin_rep = st.columns([1, 1])
            with col_admin_chat:
                st.markdown("#### 💬 会話履歴")
                for msg in target_log.get("messages", []):
                    role_name = "👤 顧客" if msg["role"] == "assistant" else "💼 銀行員"
                    st.markdown(f"> **{role_name}**\n> {msg['content']}")
                    st.write("")
                    
            with col_admin_rep:
                st.markdown("#### 📊 評価レポート")
                if target_log.get("report"):
                    st.markdown(target_log.get("report"))
                else:
                    st.info("このセッションでは評価レポートは生成されていません。")
        else:
            st.info("該当するデータがありません。")
            
    st.stop()

# --- 通常の担当者（demo）画面 ---
st.title("🏦 金融商品販売 AIロールプレイシステム")
st.markdown(f"ログイン中: **{st.session_state.username}** さん | 顧客ペルソナを設定し、銀行員としての対話・提案スキルを磨くトレーニングツールです。")

if st.sidebar.button("🚪 ログアウト", key="user_logout"):
    st.session_state.logged_in = False
    st.session_state.username = ""
    st.rerun()

# --- 6. サイドバー機能 (ペルソナ設定 & ファイルロード) ---
st.sidebar.header("👤 顧客ペルソナ設定")

p_name = st.sidebar.text_input("氏名", value="山田 規子")
p_age = st.sidebar.text_input("年齢", value="65歳")
p_job = st.sidebar.text_input("職業", value="専業主婦 (夫は定年退職)")
p_family = st.sidebar.text_input("家族構成", value="夫（68歳）と二人暮らし、独立した子供が2人")
p_purpose = st.sidebar.text_input("来店目的", value="定期預金が満期を迎えたため手続きに来店")
p_personality = st.sidebar.text_input("性格", value="慎重派で心配性。損はしたくないが、今の低金利には不満がある。")
p_experience = st.sidebar.selectbox("投資経験", ["全くない", "少しある（過去に国債のみ）", "豊富にある"])

current_persona = {
    "name": p_name, "age": p_age, "job": p_job, "family": p_family,
    "purpose": p_purpose, "personality": p_personality, "experience": p_experience
}

if st.sidebar.button("⚙️ 顧客ペルソナ設定・リセット"):
    st.session_state.persona = current_persona
    with st.spinner("新しいペルソナに合わせた最初の発言を生成中..."):
        dynamic_initial_message = generate_initial_greeting(current_persona, EMBEDDED_API_KEY)
    st.session_state.messages = [{"role": "assistant", "content": dynamic_initial_message}]
    st.session_state.last_response = ""
    st.session_state.report = ""
    st.success("ペルソナ設定に合わせて最初の発言を変更し、初期化しました！")
    st.rerun()

st.sidebar.markdown("---")
st.sidebar.header("📁 商品情報・規約ロード")

uploaded_files = st.sidebar.file_uploader(
    "金融商品資料等 (Word, PDF, PPT, Excel)", 
    type=["docx", "pdf", "pptx", "xlsx", "xls"], 
    accept_multiple_files=True,
    key="file_uploader"
)

all_extra_text = []
if uploaded_files:
    for f in uploaded_files:
        try:
            if f.name.endswith(".docx"): content = extract_from_docx(f)
            elif f.name.endswith(".pdf"): content = extract_from_pdf(f)
            elif f.name.endswith(".pptx"): content = extract_from_pptx(f)
            elif f.name.endswith((".xlsx", ".xls")): content = extract_from_excel(f)
            all_extra_text.append(f"--- ファイル名: {f.name} ---\n{content}")
            st.sidebar.write(f"✔️ 資料読込済: {f.name}")
        except Exception as e:
            st.sidebar.error(f"❌ {f.name} の読込失敗")

st.sidebar.markdown("---")
if st.sidebar.button("🛑 アプリを終了する"):
    st.sidebar.warning("システムを終了します。")
    os.getpid()
    os.kill(os.getpid(), signal.SIGINT)

# --- 7. セッション状態の初期化 ---
if "persona" not in st.session_state:
    st.session_state.persona = current_persona
if "messages" not in st.session_state:
    st.session_state.messages = [{"role": "assistant", "content": "こんにちは。定期預金の満期の件で来たんですが、今の時代、普通に預けていても全然増えないですよね…"}]
if "last_response" not in st.session_state:
    st.session_state.last_response = ""
if "report" not in st.session_state:
    st.session_state.report = ""

# --- 8. メインエリアのレイアウト設計 ---
col_chat, col_report = st.columns([1.2, 1.0])

# --- 左側：チャットエリア ---
with col_chat:
    st.subheader("💬 ロールプレイ対話画面")
    st.info(f"【現在対話中の顧客】 {st.session_state.persona['name']}さん ({st.session_state.persona['age']} / {st.session_state.persona['job']})")

    for m in st.session_state.messages:
        role = "assistant" if m["role"] == "assistant" else "user"
        avatar = "👤" if role == "assistant" else "💼"
        with st.chat_message(role, avatar=avatar):
            st.markdown(m["content"])

    if prompt := st.chat_input("銀行員としての返答・提案を入力してください"):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user", avatar="💼"):
            st.markdown(prompt)

        with st.chat_message("assistant", avatar="👤"):
            with st.spinner(f"{st.session_state.persona['name']}さんが考えています..."):
                res = get_ai_roleplay_response(
                    st.session_state.messages, 
                    st.session_state.persona, 
                    all_extra_text, 
                    EMBEDDED_API_KEY
                )
            st.markdown(res)
            
        st.session_state.messages.append({"role": "assistant", "content": res})
        st.rerun()

# --- 右側：評価結果エリア ---
with col_report:
    st.subheader("📊 応対評価レポート")
    st.markdown("ロールプレイの節目、または終了時に以下のボタンを押して評価を生成してください。")
    
    if st.button("📝 応対評価レポートを生成する", type="primary"):
        if len(st.session_state.messages) <= 1:
            st.warning("会話が開始されていません。ロールプレイを行ってから評価してください。")
        else:
            with st.spinner("これまでの会話内容からコンプライアンス・応対品質を分析中..."):
                report_res = generate_evaluation_report(
                    st.session_state.messages, 
                    st.session_state.persona, 
                    EMBEDDED_API_KEY
                )
                st.session_state.report = report_res
                st.session_state.last_response = report_res # Excel出力用
                
                # レポート生成時のみ、まとめてSupabaseに一度だけログ保存を実行
                save_log_to_supabase(
                    st.session_state.username, 
                    st.session_state.persona["name"], 
                    st.session_state.messages, 
                    report=report_res
                )
                
                st.rerun()
                
    if st.session_state.report:
        st.markdown(st.session_state.report)
        
        excel_data = create_excel_download(st.session_state.report)
        st.download_button(
            label="📥 評価結果レポートをExcelでダウンロード",
            data=excel_data,
            file_name="roleplay_evaluation.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
    else:
        st.info("ここにAIによる採点、コンプライアンスチェックのレポートが表示されます。")
